"""
美股市场深度分析 — 2026年5月19日
Investment analysis presentation on US stock market
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Color Palette: Midnight Executive ──
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF2, 0xF4, 0xF8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x4A, 0x5A, 0x6B)
ACCENT_GOLD = RGBColor(0xD4, 0xA8, 0x43)
ACCENT_RED = RGBColor(0xE8, 0x4D, 0x4D)
ACCENT_GREEN = RGBColor(0x2E, 0xAD, 0x6F)
LIGHT_NAVY = RGBColor(0x2D, 0x38, 0x7A)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
LINE_COLOR = RGBColor(0xE2, 0xE8, 0xF0)

# ── Presentation Setup ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helper Functions ──

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text="", font_size=14,
                font_color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT,
                font_name="Calibri", valign=MSO_ANCHOR.TOP, line_spacing=1.15,
                margin=Inches(0.05)):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    # Set margin
    tb.text_frame.margin_left = margin
    tb.text_frame.margin_right = margin
    tb.text_frame.margin_top = Inches(0)
    tb.text_frame.margin_bottom = Inches(0)
    if text:
        p = tb.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size * line_spacing)
    return tb

def add_rich_textbox(slide, left, top, width, height, runs_list,
                     align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                     line_spacing=1.15, margin=Inches(0.05)):
    """runs_list = list of dicts: {text, size, color, bold, font_name, break_line}"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = margin
    tf.margin_right = margin
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    p = tf.paragraphs[0]
    for i, r in enumerate(runs_list):
        run = p.add_run()
        run.text = r.get("text", "")
        run.font.size = Pt(r.get("size", 14))
        run.font.color.rgb = r.get("color", DARK_TEXT)
        run.font.bold = r.get("bold", False)
        run.font.name = r.get("font_name", "Calibri")
        p.alignment = align
        p.space_after = Pt(0)
        p.line_spacing = Pt(r.get("size", 14) * line_spacing)
        if r.get("break_line") and i < len(runs_list) - 1:
            p = tf.add_paragraph()
    return tb

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None,
             border_width=Pt(0), corner_radius=None, shadow=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    if corner_radius:
        # Adjust corner radius (default is quite large)
        pass
    if shadow:
        sp = shape._element
        spPr = sp.find(qn('a:spPr'))
        if spPr is None:
            spPr = sp.makeelement(qn('a:spPr'), {})
            sp.insert(0, spPr)
        effectLst = spPr.find(qn('a:effectLst'))
        if effectLst is None:
            effectLst = spPr.makeelement(qn('a:effectLst'), {})
            spPr.append(effectLst)
        outerShdw = spPr.makeelement(qn('a:outerShdw'), {
            'blurRad': '50800',
            'dist': '25400',
            'dir': '1350000',
            'algn': 'bl',
        })
        srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
        alpha = srgbClr.makeelement(qn('a:alpha'), {'val': '20000'})
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)
    return shape

def add_line(slide, left, top, width, color=LINE_COLOR, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), line_width
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_stat_card(slide, left, top, width, height, number, label, sublabel="",
                  number_color=NAVY, bg_color=CARD_BG):
    """Large stat callout card"""
    add_rect(slide, left, top, width, height, fill_color=bg_color, shadow=True)
    add_textbox(slide, left + 0.25, top + 0.15, width - 0.5, 1.0,
                number, font_size=40, font_color=number_color, bold=True,
                font_name="Georgia")
    add_textbox(slide, left + 0.25, top + 1.2, width - 0.5, 0.5,
                label, font_size=13, font_color=DARK_TEXT, bold=True)
    if sublabel:
        add_textbox(slide, left + 0.25, top + 1.6, width - 0.5, 0.4,
                    sublabel, font_size=10, font_color=MUTED)

def add_icon_circle(slide, left, top, size, icon_char, fill_color=NAVY):
    """Circle with icon text"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = icon_char
    run.font.size = Pt(size * 36)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    return shape

def add_card_row(slide, left, top, cards, card_w=2.5, card_h=1.8, gap=0.25):
    """Add a row of stat/icon cards"""
    for i, card in enumerate(cards):
        x = left + i * (card_w + gap)
        add_rect(slide, x, top, card_w, card_h, fill_color=CARD_BG, shadow=True)
        # Icon circle
        add_icon_circle(slide, x + 0.2, top + 0.2, 0.4, card.get("icon", "?"),
                        fill_color=card.get("icon_color", NAVY))
        add_textbox(slide, x + 0.75, top + 0.2, card_w - 1.0, 0.5,
                    card.get("title", ""), font_size=14, font_color=DARK_TEXT, bold=True)
        add_textbox(slide, x + 0.2, top + 0.8, card_w - 0.4, card_h - 1.0,
                    card.get("desc", ""), font_size=11, font_color=MUTED)

# ── SLIDE 1: Title ──
slide = add_blank_slide()
# Dark navy background
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=NAVY)
# Decorative gradient bar (simulated with shapes)
add_rect(slide, 0, 0, 0.12, 7.5, fill_color=ACCENT_GOLD)
# Title content
add_textbox(slide, 1.5, 1.5, 10.5, 1.5, "美股市场深度分析",
            font_size=52, font_color=WHITE, bold=True, font_name="Georgia")
add_line(slide, 1.5, 3.2, 3.0, color=ACCENT_GOLD, line_width=Pt(3))
add_textbox(slide, 1.5, 3.6, 10.5, 1.2,
            "地缘冲突 · 通胀升温 · 联储转向 · AI贸易审判日",
            font_size=22, font_color=ICE_BLUE, font_name="Calibri")
add_textbox(slide, 1.5, 5.3, 10.5, 0.5, "2026年5月19日  |  投资分析报告",
            font_size=14, font_color=ICE_BLUE)
add_textbox(slide, 1.5, 5.85, 10.5, 0.4, "数据来源：Bloomberg, MarketWatch, Federal Reserve, CNBC",
            font_size=10, font_color=RGBColor(0x99, 0xA8, 0xBB))

# ── SLIDE 2: Agenda ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.5, 5.0, 0.8, "本报告议题",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")
add_line(slide, 0.8, 1.25, 1.8, color=ACCENT_GOLD, line_width=Pt(2.5))

agenda_items = [
    ("01", "三大指数与市场总览", "Dow · S&P 500 · Nasdaq 关键位与近期走势"),
    ("02", "三大核心驱动力", "伊朗冲突与油价 · 通胀与联储 · AI科技股"),
    ("03", "板块深度分析", "科技 · 能源 · 金融 三大板块分化格局"),
    ("04", "债市信号与利率预期", "收益率飙升 · 加息概率 · 曲线形态"),
    ("05", "NVIDIA：AI贸易审判日", "Mag 7 分化 · 数据中心军备竞赛 · 527B投资"),
    ("06", "风险矩阵与投资策略", "集中度风险 · 情景分析 · 配置建议"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = 1.8 + i * 0.85
    add_icon_circle(slide, 0.8, y, 0.45, num, fill_color=NAVY)
    add_textbox(slide, 1.5, y - 0.02, 10.0, 0.45, title,
                font_size=18, font_color=DARK_TEXT, bold=True, font_name="Georgia")
    add_textbox(slide, 1.5, y + 0.42, 10.0, 0.35, desc,
                font_size=11, font_color=MUTED)
    if i < len(agenda_items) - 1:
        add_line(slide, 1.5, y + 0.78, 10.5, color=LINE_COLOR)

# ── SLIDE 3: 三大指数概览 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 8.0, 0.7, "三大指数概览",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")
add_textbox(slide, 0.8, 1.0, 8.0, 0.4, "2026年5月18日收盘  |  连续第二日下跌",
            font_size=14, font_color=MUTED)

# Three big stat cards
add_stat_card(slide, 0.8, 1.8, 3.6, 2.3, "49,686", "道琼斯工业指数", "+0.32%  ·  唯一上涨",
              number_color=ACCENT_GREEN)
add_stat_card(slide, 4.8, 1.8, 3.6, 2.3, "7,403", "标普500指数", "−0.07%  ·  连续第二日下跌",
              number_color=ACCENT_RED)
add_stat_card(slide, 8.8, 1.8, 3.6, 2.3, "26,091", "纳斯达克综合指数", "−0.51%  ·  半导体领跌",
              number_color=ACCENT_RED)

# Key callouts at bottom
add_textbox(slide, 0.8, 4.5, 5.5, 0.4, "关键支撑与阻力",
            font_size=14, font_color=NAVY, bold=True)
add_rich_textbox(slide, 0.8, 4.9, 5.5, 2.2, [
    {"text": "S&P 500", "size": 12, "color": DARK_TEXT, "bold": True, "break_line": False},
    {"text": "  支撑 7,250 / 阻力 7,500–7,750", "size": 12, "color": MUTED, "break_line": True},
    {"text": "Nasdaq", "size": 12, "color": DARK_TEXT, "bold": True, "break_line": False},
    {"text": "  关键位 26,000，失守则看 25,300", "size": 12, "color": MUTED, "break_line": True},
    {"text": "Dow", "size": 12, "color": DARK_TEXT, "bold": True, "break_line": False},
    {"text": "  区间 48,700–50,200，突破上看 51,300", "size": 12, "color": MUTED, "break_line": True},
    {"text": "VIX 恐慌指数", "size": 12, "color": DARK_TEXT, "bold": True, "break_line": False},
    {"text": "  17.81（−3.36%），仍处中等水平", "size": 12, "color": MUTED, "break_line": False},
])

add_rich_textbox(slide, 7.0, 4.5, 5.5, 2.5, [
    {"text": "市场概况", "size": 14, "color": NAVY, "bold": True, "break_line": True},
    {"text": "S&P 500今年已录12次历史新高，但5月中旬开始出现回调。标普指数自3月底部反弹18.4%，但等权重版仅涨8.3%——市场广度极度收窄。", "size": 12, "color": MUTED, "break_line": True},
    {"text": "科技股获利了结", "size": 12, "color": ACCENT_RED, "bold": True, "break_line": False},
    {"text": "  +  ", "size": 12, "color": MUTED, "break_line": False},
    {"text": "伊朗停火预期摇摆", "size": 12, "color": ACCENT_RED, "bold": True, "break_line": True},
    {"text": "共同压制成市场情绪。", "size": 12, "color": MUTED, "break_line": False},
])

# ── SLIDE 4: 地缘政治驱动 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 10.0, 0.7, "核心驱动力一：地缘政治与能源危机",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")
add_textbox(slide, 0.8, 1.0, 10.0, 0.4, "美国-伊朗冲突持续，霍尔木兹海峡封锁 → 油价飙升 → 通胀传导",
            font_size=14, font_color=MUTED)

# Left panel - timeline
add_rect(slide, 0.8, 1.8, 5.8, 5.2, fill_color=CARD_BG, shadow=True)
add_textbox(slide, 1.1, 2.0, 5.3, 0.4, "关键事件时间线",
            font_size=16, font_color=NAVY, bold=True, font_name="Georgia")

timeline = [
    ("2月下旬", "美国与以色列联合对伊朗发动军事行动，冲突爆发"),
    ("4月8日", "短暂停火（后破裂），伊斯兰堡谈判无果而终"),
    ("4月下旬", "霍尔木兹海峡航运几乎中断，布伦特原油触及$140+"),
    ("5月中旬", "特朗普警告「时间在流逝」，但随后叫停攻击计划"),
    ("5月18日", "WTI $108.66/bbl，布伦特 $112.10/bbl，能源板块+1.8%"),
    ("当前", "停火谈判进展不明，油价波动剧烈，市场对消息面极度敏感"),
]
for i, (date, desc) in enumerate(timeline):
    y = 2.6 + i * 0.7
    add_icon_circle(slide, 1.1, y, 0.3, "●", fill_color=NAVY)
    add_textbox(slide, 1.6, y - 0.05, 1.3, 0.35, date,
                font_size=11, font_color=NAVY, bold=True)
    add_textbox(slide, 3.0, y - 0.05, 3.3, 0.35, desc, font_size=10, font_color=MUTED)

# Right panel - impact numbers
add_rect(slide, 7.0, 1.8, 5.5, 5.2, fill_color=NAVY)
add_textbox(slide, 7.4, 2.1, 4.8, 0.5, "能源冲击量化",
            font_size=16, font_color=WHITE, bold=True, font_name="Georgia")

oil_stats = [
    ("+55%", "布伦特原油自冲突前涨幅", "2月底 ≈$70 → 当前 $112"),
    ("$4.51", "美国平均汽油价格/加仑", "汽油同比上涨 28%+"),
    ("+32.8%", "能源板块(XLE)年内涨幅", "年至今最强板块，受益于高油价"),
    ("3.8%", "4月CPI同比增幅", "三年新高，能源是主要推手"),
    ("−1.78M", "IEA预测全球供应缺口(桶/日)", "全年供不应求，OPEC+维持减产"),
]
for i, (num, label, sub) in enumerate(oil_stats):
    y = 2.8 + i * 0.85
    add_textbox(slide, 7.4, y, 4.8, 0.42, num,
                font_size=26, font_color=ACCENT_GOLD, bold=True, font_name="Georgia")
    add_textbox(slide, 7.4, y + 0.42, 4.8, 0.22, label,
                font_size=12, font_color=WHITE, bold=True)
    add_textbox(slide, 7.4, y + 0.63, 4.8, 0.2, sub,
                font_size=9, font_color=RGBColor(0x99, 0xA8, 0xBB))

# ── SLIDE 5: 通胀与美联储 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 10.0, 0.7, "核心驱动力二：通胀飙升与美联储困境",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")
add_textbox(slide, 0.8, 1.0, 10.0, 0.4,
            "联邦基金利率 3.50–3.75%  |  市场从「期待降息」转向「定价加息」",
            font_size=14, font_color=MUTED)

# Inflation stat callouts
inflation_data = [
    ("3.8%", "4月CPI\n同比", ACCENT_RED),
    ("6.0%", "4月PPI\n同比", ACCENT_RED),
    ("3.2%", "核心PCE\n月环比", ACCENT_RED),
    ("37%", "年底加息\n概率", ACCENT_GOLD),
]
for i, (num, label, color) in enumerate(inflation_data):
    x = 0.8 + i * 3.1
    add_rect(slide, x, 1.8, 2.8, 1.8, fill_color=CARD_BG, shadow=True)
    add_textbox(slide, x + 0.2, 1.9, 2.4, 0.9, num,
                font_size=40, font_color=color, bold=True, font_name="Georgia")
    add_textbox(slide, x + 0.2, 2.8, 2.4, 0.6, label,
                font_size=12, font_color=MUTED)

# Fed details - left
add_rect(slide, 0.8, 4.0, 5.8, 3.0, fill_color=CARD_BG, shadow=True)
add_textbox(slide, 1.1, 4.15, 5.3, 0.4, "美联储政策动态",
            font_size=16, font_color=NAVY, bold=True, font_name="Georgia")
fed_details = [
    "4月FOMC以 8-4 投票维持利率不变（1992年以来最多反对票）",
    "3位鹰派地区行长反对声明中的「宽松倾向」",
    "Kevin Warsh 被提名接替 Powell 任新主席（5月11日听证会）",
    "Warsh 被视为偏鸽派 → 政策不确定性上升",
    "BofA 将首次降息预测推迟至 2027下半年",
    "Jeffrey Gundlach：「通胀数据让降息变得不可能」",
]
for i, d in enumerate(fed_details):
    add_textbox(slide, 1.1, 4.7 + i * 0.35, 5.3, 0.3,
                f"•  {d}", font_size=10, font_color=DARK_TEXT)

# Rate outlook - right
add_rect(slide, 7.0, 4.0, 5.5, 3.0, fill_color=NAVY)
add_textbox(slide, 7.3, 4.15, 4.8, 0.4, "利率预期剧变",
            font_size=16, font_color=WHITE, bold=True, font_name="Georgia")
rate_data = [
    ("30-50%", "市场定价年底前至少加息一次"),
    ("~47%", "Kalshi预测市场：累计加息概率（至2027年7月）"),
    ("4.55%", "10年期美债收益率——一年新高"),
    ("5%+", "30年期美债收益率——二十年高位附近"),
    ("3.88%", "2年期收益率——已高于联邦基金利率"),
]
for i, (num, desc) in enumerate(rate_data):
    y = 4.7 + i * 0.45
    add_textbox(slide, 7.3, y, 1.5, 0.35, num,
                font_size=18, font_color=ACCENT_GOLD, bold=True, font_name="Georgia")
    add_textbox(slide, 8.9, y + 0.02, 3.4, 0.35, desc,
                font_size=11, font_color=ICE_BLUE)

# ── SLIDE 6: AI与科技股 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 10.0, 0.7, "核心驱动力三：AI科技股与Mag 7分化",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")
add_textbox(slide, 0.8, 1.0, 10.0, 0.4,
            "AI贸易从「买入一切」转向「甄别赢家」——市场奖励变现能力，惩罚无限capex",
            font_size=14, font_color=MUTED)

# Mag 7 scorecard
mag7 = [
    ("Alphabet", "+9.96%", "AI变现赢家", ACCENT_GREEN, "Google"),
    ("Apple", "+5.00%", "iPhone周期强劲", ACCENT_GREEN, "Apple"),
    ("Amazon", "+0.77%", "AWS强但capex担忧", DARK_TEXT, "Amazon"),
    ("Microsoft", "−3.93%", "AI支出焦虑", ACCENT_RED, "MSFT"),
    ("Meta", "−8.55%", "capex ROI重置", ACCENT_RED, "Meta"),
    ("Tesla", "−3.56%", "叙事最不确定", ACCENT_RED, "Tesla"),
    ("NVIDIA", "5/20公布", "AI审判日", ACCENT_GOLD, "NVIDIA"),
]
for i, (name, perf, note, color, short) in enumerate(mag7):
    x = 0.6 + i * 1.78
    add_rect(slide, x, 1.7, 1.65, 2.6, fill_color=CARD_BG, shadow=True)
    add_textbox(slide, x + 0.1, 1.8, 1.45, 0.3, short,
                font_size=11, font_color=MUTED, bold=True)
    add_textbox(slide, x + 0.1, 2.15, 1.45, 0.5, perf,
                font_size=24, font_color=color, bold=True, font_name="Georgia")
    add_textbox(slide, x + 0.1, 2.7, 1.45, 0.35, name,
                font_size=10, font_color=DARK_TEXT, bold=True)
    add_textbox(slide, x + 0.1, 3.0, 1.45, 0.35, note,
                font_size=9, font_color=MUTED)

# AI capex context
add_rect(slide, 0.8, 4.7, 5.8, 2.3, fill_color=NAVY)
add_textbox(slide, 1.1, 4.85, 5.3, 0.35, "AI数据中心军备竞赛",
            font_size=16, font_color=WHITE, bold=True, font_name="Georgia")
add_textbox(slide, 1.1, 5.25, 5.3, 0.75, "$5270亿",
            font_size=42, font_color=ACCENT_GOLD, bold=True, font_name="Georgia")
add_textbox(slide, 1.1, 6.0, 5.3, 0.35,
            "Mag 7 2026财年AI+数据中心总资本支出预估",
            font_size=12, font_color=ICE_BLUE)
add_textbox(slide, 1.1, 6.35, 5.3, 0.45,
            "超过全球30个国家的GDP  |  2025-2030年数据中心容量预计翻倍（+97GW）",
            font_size=10, font_color=RGBColor(0x99, 0xA8, 0xBB))

add_rect(slide, 7.0, 4.7, 5.5, 2.3, fill_color=CARD_BG, shadow=True)
add_textbox(slide, 7.3, 4.85, 4.8, 0.4, "关键结论",
            font_size=16, font_color=NAVY, bold=True, font_name="Georgia")
conclusions = [
    "市场从「买AI故事」转向「买AI利润」——分化加剧",
    "Alphabet + Apple 通过变现测试；Meta + MSFT 遭遇capex质疑",
    "Mag 7占S&P 500权重 34.8%（2016年仅12.5%）",
    "NVIDIA 5月20日财报 = AI基础设施需求的终极验证",
    "PEG仅0.68 → 即使经历历史性涨幅，估值仍有吸引力",
]
for i, c in enumerate(conclusions):
    add_textbox(slide, 7.3, 5.35 + i * 0.33, 4.8, 0.28,
                f"•  {c}", font_size=10, font_color=DARK_TEXT)

# ── SLIDE 7: 板块表现 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 10.0, 0.7, "板块深度分析",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")
add_textbox(slide, 0.8, 1.0, 10.0, 0.4,
            "三大关键板块呈现极致分化：科技领跑、能源受益、金融承压",
            font_size=14, font_color=MUTED)

# Three column analysis
sector_data = [
    {
        "name": "科技 (AI/半导体)",
        "perf": "+40.7% YTD",
        "perf_detail": "SMH半导体ETF年内涨幅",
        "color": ACCENT_GREEN,
        "items": [
            "AI基础设施超级周期驱动，SMH ETF +40.7% YTD",
            "科技板块5月第一周单周涨 +8.5%",
            "65%科技公司超预期盈利",
            "Mag 7 Q1盈利增长 22.8%（其余S&P 500仅10.1%）",
            "⚠️ 费城半导体指数5/18暴跌 2.5-3.3%",
            "⚠️ 仅~10只股票贡献S&P 500约70%涨幅",
        ]
    },
    {
        "name": "能源",
        "perf": "+32.8% YTD",
        "perf_detail": "XLE能源ETF年内涨幅",
        "color": ACCENT_GOLD,
        "items": [
            "受益于伊朗冲突+霍尔木兹海峡封锁",
            "全球库存低于历史均值，OPEC+维持减产",
            "IEA预测全年178万桶/日供应缺口",
            "Chevron(CVX) Q1盈利超预期 53.3%",
            "⚠️ 5月首周回调 -5.3%，油价波动剧烈",
            "⚠️ 停火谈判进展可令油价骤跌 → 板块急转",
        ]
    },
    {
        "name": "金融",
        "perf": "−1.3% (月初)",
        "perf_detail": "5月首周表现",
        "color": ACCENT_RED,
        "items": [
            "FOMC政策不确定性压制板块估值",
            "银行 −3.1%，保险 −1.6%",
            "4月FOMC 8-4票（最多反对票）动摇信心",
            "加息预期升温 → 信贷成本上升",
            "✅ DBS仍超配：收入韧性+健康资产负债表",
            "⚠️ 联储主席交接期 → 政策方向不明朗",
        ]
    },
]

for i, sec in enumerate(sector_data):
    x = 0.6 + i * 4.2
    # Header
    add_rect(slide, x, 1.7, 3.9, 0.9, fill_color=NAVY)
    add_textbox(slide, x + 0.2, 1.75, 3.5, 0.35, sec["name"],
                font_size=18, font_color=WHITE, bold=True, font_name="Georgia")
    add_textbox(slide, x + 0.2, 2.1, 3.5, 0.35, sec["perf"],
                font_size=22, font_color=sec["color"], bold=True, font_name="Georgia")
    # Items
    add_rect(slide, x, 2.6, 3.9, 4.2, fill_color=CARD_BG, shadow=True)
    for j, item in enumerate(sec["items"]):
        add_textbox(slide, x + 0.2, 2.8 + j * 0.55, 3.5, 0.5,
                    f"•  {item}", font_size=10, font_color=DARK_TEXT)

# ── SLIDE 8: 债市信号 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 10.0, 0.7, "债券市场：加息信号的报警器",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")
add_textbox(slide, 0.8, 1.0, 10.0, 0.4,
            "美债收益率全面飙升——「债券卫士」正在用脚投票",
            font_size=14, font_color=MUTED)

# Yield cards
yield_data = [
    ("4.55%", "10年期美债", "一年新高，周涨24bp", ACCENT_RED),
    ("5%+", "30年期美债", "二十年高位附近", ACCENT_RED),
    ("3.88%", "2年期美债", "已高于联邦基金利率", ACCENT_GOLD),
    ("5%以上", "\"恐慌按钮\"", "市场共识关键心理关口", ACCENT_RED),
]
for i, (num, label, sub, color) in enumerate(yield_data):
    x = 0.8 + i * 3.1
    add_rect(slide, x, 1.7, 2.8, 1.8, fill_color=CARD_BG, shadow=True)
    add_textbox(slide, x + 0.2, 1.85, 2.4, 0.7, num,
                font_size=38, font_color=color, bold=True, font_name="Georgia")
    add_textbox(slide, x + 0.2, 2.55, 2.4, 0.35, label,
                font_size=14, font_color=DARK_TEXT, bold=True)
    add_textbox(slide, x + 0.2, 2.9, 2.4, 0.3, sub,
                font_size=10, font_color=MUTED)

# Bear steepening explanation
add_rect(slide, 0.8, 3.9, 5.8, 3.1, fill_color=CARD_BG, shadow=True)
add_textbox(slide, 1.1, 4.05, 5.3, 0.4, "熊陡 (Bear Steepening) 的含义",
            font_size=16, font_color=NAVY, bold=True, font_name="Georgia")
bear_points = [
    "长端收益率上升速度快于短端 → 曲线陡峭化",
    "核心驱动：通胀预期升温 + 财政赤字担忧",
    "经典信号：市场预判联储将被迫加息",
    "结果：远期现金流折现率上升 → 高估值成长股承压",
    "5月15日：10年期美债单日飙升 73bp（罕见极端波动）",
    "日本30年期国债同步突破 4%（1999年以来首次）",
]
for i, p in enumerate(bear_points):
    add_textbox(slide, 1.1, 4.6 + i * 0.38, 5.3, 0.32,
                f"•  {p}", font_size=11, font_color=DARK_TEXT)

# Impact analysis - right
add_rect(slide, 7.0, 3.9, 5.5, 3.1, fill_color=NAVY)
add_textbox(slide, 7.3, 4.05, 4.8, 0.4, "对股市的传导机制",
            font_size=16, font_color=WHITE, bold=True, font_name="Georgia")
impact_data = [
    ("长久期科技股", "折现率↑ → 远期现金流估值↓ → 股价承压"),
    ("REITs", "Cap rate扩张 + 浮息债务成本↑ → 双重打击"),
    ("小盘股(Russell 2000)", "融资成本敏感度高，日内波动极端"),
    ("美元指数 +1.5%", "强美元 → 跨国公司海外收入汇兑损失"),
    ("避险资产受益", "防御性板块（必需消费 +9%）获资金流入"),
]
for i, (title, desc) in enumerate(impact_data):
    y = 4.6 + i * 0.5
    add_textbox(slide, 7.3, y, 1.8, 0.35, title,
                font_size=11, font_color=ACCENT_GOLD, bold=True)
    add_textbox(slide, 9.2, y, 3.1, 0.35, desc,
                font_size=10, font_color=ICE_BLUE)

# ── SLIDE 9: NVIDIA深度 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=NAVY)
add_rect(slide, 0, 0, 0.12, 7.5, fill_color=ACCENT_GOLD)
add_textbox(slide, 1.2, 0.5, 10.0, 0.7, "焦点：NVIDIA Q1 FY2027 财报",
            font_size=40, font_color=WHITE, bold=True, font_name="Georgia")
add_textbox(slide, 1.2, 1.15, 10.0, 0.4, "5月20日公布  |  美股AI贸易的「审判日」",
            font_size=16, font_color=ICE_BLUE)

# Key metrics
nv_metrics = [
    ("$78-79B", "预计营收", "上季$68.13B"),
    ("$1.74-1.78", "预计调整后EPS", "毛利率~74.9%"),
    ("$320", "BofA目标价", "距当前~$226有41.6%空间"),
    ("$527B", "Mag 7年度AI capex", "历史级别投资周期"),
]
for i, (num, label, sub) in enumerate(nv_metrics):
    x = 1.2 + i * 2.95
    add_textbox(slide, x, 2.0, 2.5, 0.7, num,
                font_size=36, font_color=ACCENT_GOLD, bold=True, font_name="Georgia")
    add_textbox(slide, x, 2.7, 2.5, 0.3, label,
                font_size=13, font_color=WHITE, bold=True)
    add_textbox(slide, x, 3.0, 2.5, 0.25, sub,
                font_size=10, font_color=RGBColor(0x99, 0xA8, 0xBB))

# Analyst targets
add_textbox(slide, 1.2, 3.7, 5.0, 0.4, "华尔街一致看多",
            font_size=18, font_color=WHITE, bold=True, font_name="Georgia")
firms = [
    ("BofA", "$320", "Buy"), ("KeyBanc", "$300", "Overweight"),
    ("Morgan Stanley", "$285", "Overweight"), ("TD Cowen", "$275", "Buy"),
    ("UBS", "$275", "Buy"), ("37分析师共识", "$274", "Strong Buy"),
]
for i, (firm, target, rating) in enumerate(firms):
    x = 1.2 + (i % 3) * 3.9
    y = 4.2 + (i // 3) * 0.55
    add_textbox(slide, x, y, 1.8, 0.35, firm,
                font_size=12, font_color=ICE_BLUE, bold=True)
    add_textbox(slide, x + 1.9, y, 0.9, 0.35, target,
                font_size=12, font_color=ACCENT_GOLD, bold=True)
    add_textbox(slide, x + 2.9, y, 0.9, 0.35, rating,
                font_size=11, font_color=ACCENT_GREEN)

# Huang quote
add_line(slide, 1.2, 5.3, 11.0, color=RGBColor(0x3D, 0x4A, 0x8C), line_width=Pt(1))
add_rich_textbox(slide, 1.2, 5.5, 11.0, 1.0, [
    {"text": '「如果我们交了一个糟糕的季度，就是AI泡沫的证据。如果交了一个伟大的季度，我们就在助长AI泡沫。」', "size": 16, "color": ICE_BLUE, "bold": False, "font_name": "Georgia", "break_line": True},
    {"text": "— Jensen Huang, NVIDIA CEO", "size": 12, "color": MUTED},
])

# PEG context
add_textbox(slide, 1.2, 6.6, 11.0, 0.4,
            "当前PEG仅 0.68（低于1.0公允价值阈值）→ 市场可能仍低估其增长潜力",
            font_size=13, font_color=RGBColor(0x99, 0xA8, 0xBB))

# ── SLIDE 10: 市场集中度风险 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 10.0, 0.7, "市场集中度：隐形风险",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")
add_textbox(slide, 0.8, 1.0, 10.0, 0.4,
            "当市值加权指数和等权重指数背离过大，就是系统性风险信号",
            font_size=14, font_color=MUTED)

# Visual comparison - two big cards
add_rect(slide, 0.8, 1.7, 5.8, 2.5, fill_color=CARD_BG, shadow=True)
add_textbox(slide, 1.1, 1.85, 5.3, 0.35, "市值加权 S&P 500 自3/30反弹",
            font_size=14, font_color=MUTED, bold=True)
add_textbox(slide, 1.1, 2.3, 5.3, 1.0, "+18.4%",
            font_size=60, font_color=ACCENT_GREEN, bold=True, font_name="Georgia")
add_textbox(slide, 1.1, 3.4, 5.3, 0.5,
            "仅~10只股票贡献约70%涨幅  ·  表面繁荣",
            font_size=12, font_color=DARK_TEXT)

add_rect(slide, 7.0, 1.7, 5.5, 2.5, fill_color=CARD_BG, shadow=True)
add_textbox(slide, 7.3, 1.85, 4.8, 0.35, "等权重 S&P 500 自3/30反弹",
            font_size=14, font_color=MUTED, bold=True)
add_textbox(slide, 7.3, 2.3, 4.8, 1.0, "+8.3%",
            font_size=60, font_color=ACCENT_RED, bold=True, font_name="Georgia")
add_textbox(slide, 7.3, 3.4, 4.8, 0.5,
            "广泛市场参与度极低  ·  真实的疲弱",
            font_size=12, font_color=DARK_TEXT)

# Context
add_rect(slide, 0.8, 4.5, 11.7, 2.5, fill_color=CARD_BG, shadow=True)
add_textbox(slide, 1.1, 4.65, 11.3, 0.4, "历史警示",
            font_size=16, font_color=NAVY, bold=True, font_name="Georgia")

warnings = [
    ("Mag 7 占 S&P 500 权重", "34.8%（2016年仅12.5%）→ 集中度创历史记录"),
    ("AI capex循环的脆弱性", "若NVIDIA/Hyperscalers减速 → 连锁抛售 → 指数级回撤"),
    ("被动投资放大风险", "大量资金通过ETF被动配置 → 少数大市值股支配整体指数方向"),
    ("2000年互联网泡沫教训", "极端集中度 + 高估值 + 宏观收紧 = 潜在系统性回调"),
    ("等权重 vs 市值加权背离", "18.4% vs 8.3% 的差距 → 1987年以来最大分化之一"),
]
for i, (title, desc) in enumerate(warnings):
    y = 5.15 + i * 0.35
    add_textbox(slide, 1.1, y, 3.0, 0.3, title,
                font_size=11, font_color=NAVY, bold=True)
    add_textbox(slide, 4.2, y, 8.0, 0.3, desc,
                font_size=11, font_color=DARK_TEXT)

# ── SLIDE 11: 风险矩阵 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 10.0, 0.7, "关键风险矩阵",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")

risks = [
    ("地缘政治", "高", "伊朗冲突升级 / 停火谈判破裂 / 霍尔木兹持续封锁 → 油价进一步飙升", ACCENT_RED),
    ("通胀失控", "高", "5月CPI若破4 → 美联储加息概率跃升至50%+ → 全面风险资产重定价", ACCENT_RED),
    ("AI capex逆转", "中高", "NVIDIA/Hyperscalers若暗示资本支出见顶 → 半导体板块暴跌 → 指数承压", ACCENT_RED),
    ("联储政策失误", "中", "加息过迟 → 通胀根深蒂固；加息过早过猛 → 经济硬着陆", ACCENT_GOLD),
    ("市场集中度", "中", "Mag 7权重历史最高 → 任何一只暴雷都将通过ETF波及全市场", ACCENT_GOLD),
    ("美元走强", "中", "美元指数+1.5% → 跨国公司盈利受损 → S&P 500海外收入约40%", ACCENT_GOLD),
    ("日本国债传导", "新增", "日本30年国债破4%（1999来首次）→ 全球利率上行压力 → 资金流向逆转", ACCENT_GOLD),
]

header_y = 1.5
add_rect(slide, 0.8, header_y, 11.7, 0.45, fill_color=NAVY)
for j, (h, x) in enumerate([("风险因素", 0.8), ("等级", 4.0), ("触发条件与影响", 6.0)]):
    add_textbox(slide, x + 0.15, header_y + 0.05, 2.0, 0.35, h,
                font_size=11, font_color=WHITE, bold=True)

for i, (factor, level, desc, color) in enumerate(risks):
    y = 2.05 + i * 0.72
    bg = CARD_BG if i % 2 == 0 else OFF_WHITE
    add_rect(slide, 0.8, y, 11.7, 0.62, fill_color=bg, shadow=True)
    add_textbox(slide, 1.0, y + 0.12, 2.5, 0.35, factor,
                font_size=13, font_color=DARK_TEXT, bold=True)
    add_textbox(slide, 4.2, y + 0.12, 1.5, 0.35, level,
                font_size=13, font_color=color, bold=True, font_name="Georgia")
    add_textbox(slide, 6.2, y + 0.12, 6.0, 0.35, desc,
                font_size=10, font_color=MUTED)

# ── SLIDE 12: 情景分析 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 10.0, 0.7, "未来展望：三种情景",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")

scenarios = [
    {
        "name": "🟢 乐观情景 (35%)",
        "desc": "伊朗停火 + AI盈利兑现",
        "detail": "布伦特回落至$90 → CPI降温至3.5%以下 → 联储按兵不动 → NVIDIA财报超预期 → 科技板块引领新一轮反弹",
        "target": "S&P 500: 7,750–8,000",
        "color": ACCENT_GREEN,
    },
    {
        "name": "🟡 基本情景 (45%)",
        "desc": "僵局持续 + 板块轮动",
        "detail": "油价维持$100-110 → 通胀高位徘徊 → 联储观望至年底 → 资金从高估值成长股向防御性板块轮动 → 指数横盘震荡",
        "target": "S&P 500: 7,100–7,500",
        "color": ACCENT_GOLD,
    },
    {
        "name": "🔴 悲观情景 (20%)",
        "desc": "全面冲突升级 + 滞胀",
        "detail": "伊朗冲突扩大 → 布伦特突破$140 → CPI破5% → 联储被迫加息 → AI capex收缩 → 科技股泡沫破裂 → 美债遭抛售",
        "target": "S&P 500: 6,000–6,500",
        "color": ACCENT_RED,
    },
]

for i, sc in enumerate(scenarios):
    y = 1.5 + i * 1.88
    add_rect(slide, 0.8, y, 11.7, 1.65, fill_color=CARD_BG, shadow=True)
    # Color bar on left
    add_rect(slide, 0.8, y, 0.08, 1.65, fill_color=sc["color"])
    add_textbox(slide, 1.2, y + 0.08, 5.0, 0.35, sc["name"],
                font_size=18, font_color=sc["color"], bold=True, font_name="Georgia")
    add_textbox(slide, 6.5, y + 0.1, 5.5, 0.3, sc["target"],
                font_size=14, font_color=NAVY, bold=True, font_name="Georgia")
    add_textbox(slide, 1.2, y + 0.5, 11.0, 0.3, f"场景：{sc['desc']}",
                font_size=12, font_color=DARK_TEXT, bold=True)
    add_textbox(slide, 1.2, y + 0.85, 11.0, 0.65, sc["detail"],
                font_size=10, font_color=MUTED)

# ── SLIDE 13: 投资策略建议 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=OFF_WHITE)
add_textbox(slide, 0.8, 0.4, 10.0, 0.7, "投资策略建议",
            font_size=38, font_color=NAVY, bold=True, font_name="Georgia")
add_textbox(slide, 0.8, 1.0, 10.0, 0.4,
            "当前环境：高通胀 + 地缘风险 + 加息预期 + AI分化  →  防御为主、精选进攻",
            font_size=14, font_color=MUTED)

strategies = [
    ("超配能源", "积极", "伊朗冲突持续 → 油价维持高位 → XOM、CVX等受益于供应紧缺。但需设定止损，一旦停火立即减仓。", ACCENT_GREEN),
    ("精选AI龙头", "中性偏多", "NVIDIA、Alphabet等有真实变现能力的标的，避开纯capex故事。NVIDIA PEG 0.68仍有估值支撑。等5/20财报后再加仓。", ACCENT_GREEN),
    ("增加防御配置", "推荐", "必需消费(XLP +9% YTD)、医疗保健——在高通胀+加息环境中具备定价权和稳定现金流。", ACCENT_GOLD),
    ("减仓长久期资产", "建议", "REITs、未盈利科技股、小盘股——对利率最敏感，在加息预期升温时首当其冲。", ACCENT_RED),
    ("对冲尾部风险", "推荐", "配置黄金($4,580/oz)、VIX看涨期权、或做多波动率策略。地缘风险+联储不确定性 = 尾部事件概率升高。", ACCENT_GOLD),
    ("关注日本溢出效应", "监控", "日本30年期国债破4%——若继续上行，可能引发全球债券市场连锁反应。密切跟踪JGB收益率。", ACCENT_GOLD),
]

for i, (title, action, desc, color) in enumerate(strategies):
    y = 1.7 + i * 0.88
    add_rect(slide, 0.8, y, 11.7, 0.78, fill_color=CARD_BG, shadow=True)
    add_textbox(slide, 1.1, y + 0.08, 2.2, 0.3, title,
                font_size=14, font_color=NAVY, bold=True)
    add_textbox(slide, 3.4, y + 0.08, 1.2, 0.3, action,
                font_size=12, font_color=color, bold=True)
    add_textbox(slide, 4.8, y + 0.08, 7.4, 0.55, desc,
                font_size=10, font_color=MUTED)

# ── SLIDE 14: 总结 ──
slide = add_blank_slide()
add_rect(slide, 0, 0, 13.333, 7.5, fill_color=NAVY)
add_rect(slide, 0, 0, 0.12, 7.5, fill_color=ACCENT_GOLD)
add_textbox(slide, 1.2, 0.8, 10.0, 0.7, "核心结论",
            font_size=44, font_color=WHITE, bold=True, font_name="Georgia")
add_line(slide, 1.2, 1.65, 2.5, color=ACCENT_GOLD, line_width=Pt(3))

conclusions_final = [
    ("市场处于十字路口", "连续两日下跌不是偶然——通胀、利率、地缘三大逆风正在汇聚"),
    ("分化取代普涨", "AI内部、板块之间、市值加权vs等权重——2026年选股能力比方向判断更重要"),
    ("NVIDIA是关键催化剂", "5月20日财报将决定AI贸易是延续还是逆转——做好双向准备"),
    ("不要忽视债券信号", "10年期收益率逼近5% + 曲线熊陡 → 这是市场在告诉我们：通胀没有结束"),
    ("保持防御姿态", "增加能源和必需消费敞口，减仓对利率敏感的资产，持有现金等待更好的入场时机"),
]

for i, (title, desc) in enumerate(conclusions_final):
    y = 2.1 + i * 0.88
    add_icon_circle(slide, 1.2, y, 0.45, str(i + 1), fill_color=ACCENT_GOLD)
    add_textbox(slide, 1.9, y - 0.05, 10.0, 0.4, title,
                font_size=20, font_color=WHITE, bold=True, font_name="Georgia")
    add_textbox(slide, 1.9, y + 0.38, 10.0, 0.35, desc,
                font_size=12, font_color=ICE_BLUE)

add_line(slide, 1.2, 6.35, 11.0, color=RGBColor(0x3D, 0x4A, 0x8C), line_width=Pt(0.5))
add_textbox(slide, 1.2, 6.5, 10.0, 0.5, "Thank You  ·  数据截至2026年5月19日  ·  不构成投资建议",
            font_size=12, font_color=ICE_BLUE, align=PP_ALIGN.CENTER)

# ── Save ──
output_path = "e:/me/Python/Agent/美股分析_20260519.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
